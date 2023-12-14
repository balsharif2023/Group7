from flask import Flask, jsonify, render_template, request, session
import mysql.connector
import API_Connect
import PyPDF2
import json
import docx2txt
import smtplib
import os
from dotenv import load_dotenv
from pptx import Presentation
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart

# load variables from .env file
load_dotenv()

db_endpoint = os.getenv('DB_ENDPOINT')
db_username = os.getenv('DB_USERNAME')
db_password = os.getenv('DB_PASSWORD')
db_name = os.getenv('DB_NAME')

try:
    endpoint = db_endpoint
    username = db_username
    password_ = db_password
    data_base = db_name

    Connection = mysql.connector.connect(
        user=username,
        password=password_,
        host=endpoint,
        database=data_base
    )
    Cursor = Connection.cursor()
except:
    pass
    # print("Database Connection failed")
    # raise Exception("Connection to database failed")

application = Flask(__name__)
application.secret_key = os.getenv('SECRET_KEY')


@application.route('/data')
def test_message():
    return jsonify({"message": "Welcome to SmartStudyAI"})


# Routes Todo
@application.route('/api/login', methods=['POST'])
def login():
    # Login a user
    session.pop('ID', None)
    var = request.get_json()

    Login_Sql = "SELECT * FROM Users WHERE email = '%s' " % (var['addEmail'])
    try:
        Cursor.execute(Login_Sql)
    except:
        print("Database Query failed")
        return ({"Status": "Database Connection failed"})
    results = Cursor.fetchall()
    print(results, "this")

    if results == []:
        print("Invalid credentials")
        return ({"Status": "Invalid credentials"})

    else:
        if results[0][3] == var['password']:
            print("Logged in")
            session['ID'] = results[0][0]
            return jsonify({"Status": "Logged In", "ID": results[0][0]})
        else:
            print("Login Failed")
            return ({"Status": "Login Failed"})


@application.route('/api/signup', methods=['POST'])
def signup():
    # Signup a user
    make = request.get_json()

    print(make['createUsername'])
    create_SQL = "INSERT INTO Users (username, email, password) Values ('%s', '%s', '%s')" \
        % (make['createUsername'], make['createEmail'], make['createPassword'])

    try:
        Cursor.execute(create_SQL)
        Connection.commit()
    except:
        print("Database Query failed")
        return ({"Status": "Database Connection failed"})

    user_id = Cursor.lastrowid
    fetch_SQL = "SELECT * FROM Users WHERE user_id = '%s'" % user_id

    try:
        Cursor.execute(fetch_SQL)
        results = Cursor.fetchall()
    except:
        print("Database Query failed")
        return ({"Status": "Database Connection failed"})

    if results == []:
        print("Invalid credentials")
        return ({"Status": "Invalid credentials"})

    else:
        if results[0][3] == make['createPassword']:
            print("Logged in")
            session['ID'] = results[0][0]
            return jsonify({"Status": "Logged In", "ID": results[0][0]})
        else:
            print("Login Failed")
            return ({"Status": "Login Failed"})

    # return ("User Created")


@application.route('/api/documents', methods=['GET'])
def get_documents():
    # Get all documents for a user
    # Get = request.get_json()
    get_SQL = "SELECT document_id, title, subtitle, document_text FROM Documents WHERE user_id = '%s' " % session[
        'ID']
    try:
        Cursor.execute(get_SQL)
        results = Cursor.fetchall()
    except:
        print("Database Query failed")
        return ("Database Connection failed")

    return (results)


@application.route('/api/documents', methods=['POST'])
def create_document():
    # Create a new document
    # Create = request.get_json()
    try:
        title = request.form['title']
        subtitle = request.form['subtitle']
        uploaded_file = request.files['file']
    except:
        print("Frontend request failed")
        return ("Frontend request failed")

    try:
        Doc_text = ''
        if uploaded_file.filename != '':
            file_extension = uploaded_file.filename.rsplit('.', 1)[1].lower()

            if file_extension == 'pdf':
                reader = PyPDF2.PdfReader(uploaded_file)
                for pages in range(len(reader.pages)):
                    Doc_text += reader.pages[pages].extract_text()
            elif file_extension in ('docx'):
                if file_extension == 'docx':
                    Doc_text = docx2txt.process(uploaded_file)
            elif file_extension in ('pptx'):
                presentation = Presentation(uploaded_file)
                Doc_text = ""
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text_frame"):
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    Doc_text += run.text + "\n"
            else:
                print(
                    "Unsupported file format. Please upload a PDF, DOCX, or PPTX file.")
                return "Unsupported file format. Please upload a PDF, DOCX, or PPTX file."

    except:
        print("File translator failed")
        return ("File translator failed")

    create_SQL = "INSERT INTO Documents (user_id, title, subtitle, document_text) Values (%s,%s,%s,%s)"
    create_SQL_data = (session['ID'], title, subtitle, Doc_text)

    try:
        Cursor.execute(create_SQL, create_SQL_data)
        Connection.commit()
    except:
        print("Database Query failed")
        return ("Database Connection failed")

    return ("Doc created")


@application.route('/api/documents/<int:document_id>', methods=['GET'])
def get_document(document_id):
    # Get a document
    doc_SQL = "SELECT title, document_text FROM Documents WHERE document_id = '%s' " % document_id

    try:
        Cursor.execute(doc_SQL)
        results = Cursor.fetchall()
    except:
        print("Database Query failed")
        return ("Database Connection failed")

    return (results)


@application.route('/api/documents/<int:document_id>', methods=['PUT'])
def update_document(document_id):
    # Update a document
    New_Text = "New_Text"
    update_SQL = "UPDATE Documents SET document_text = '%s' WHERE document_id = '%s' " % (
        New_Text, document_id)
    try:
        Cursor.execute(update_SQL)
        Connection.commit()
    except:
        print("Database Query failed")
        return ("Database Connection failed")

    return ("Doc updated")


@application.route('/api/documents/<int:document_id>', methods=['DELETE'])
def delete_document(document_id):
    # Delete a document
    delete_SQL = "DELETE FROM Documents WHERE document_id = '%s'" % (
        document_id)

    try:
        Cursor.execute(delete_SQL)
        Connection.commit()
    except:
        print("Database Query failed")
        return ("Database Connection failed")

    return ("Doc deleted")


@application.route('/api/documents/<int:document_id>/activity/<string:activity_type>', methods=['GET'])
def generate_questions(document_id, activity_type):
    # Generate questions for a document
    query = "SELECT document_text FROM Documents WHERE document_id = '%s'" % document_id
    try:
        Cursor.execute(query)
        prompt = Cursor.fetchall()
    except:
        print("Database Query failed")
        return ("Database Connection failed")

    try:
        if activity_type == 'multiple-choice':
            response = API_Connect.Generate(prompt[0][0])
        elif activity_type == 'true-false':
            response = API_Connect.True_False(prompt[0][0])
        elif activity_type == 'fill-in-the-blank':
            response = API_Connect.Fill_Blank(prompt[0][0])
        elif activity_type == 'matching':
            response = API_Connect.Matching(prompt[0][0])
        elif activity_type == 'flash-cards':
            response = API_Connect.Text_Quest(prompt[0][0])
        elif activity_type == 'essay':
            response = API_Connect.Essay(prompt[0][0])
    except:
        print("API Function Failed")
        return ("API Function Failed")

    return (response)


@application.route('/api/documents/<int:document_id>/activity/<string:activity_type>/results', methods=['POST'])
def save_results(document_id, activity_type):
    submission = request.get_json()
    submission_form = submission['formData']
    submission_score = submission['score']
    print(submission)

    # Start submission instance and store its ID for use in session
    submission_query = "INSERT INTO Submissions (user_id, document_id, score) Values (%s,%s,%s)" % (
        session['ID'], document_id, submission_score)
    try:
        Cursor.execute(submission_query)
        Connection.commit()
    except:
        print("Database Query Failed: Submissions")
        return ("Database Connection Failed")

    try:
        submission_id = Cursor.lastrowid
        session.pop('Submission', None)
        session['Submission'] = submission_id
    except Exception as e:
        print("Database Query Failed: Submissions ID Failed -", e)
        return ("Database Connection failed")

    for i in submission_form:
        question_query = "INSERT INTO Questions (submission_id, question, question_type) VALUES (%s, %s, %s)"
        data = (session['Submission'], i['question'], activity_type)
        Cursor.execute(question_query, data)
        Connection.commit()

    get_ques_ids = "SELECT question_id FROM Questions WHERE submission_id = %s " % session[
        'Submission']
    Cursor.execute(get_ques_ids)
    Ques_ids = Cursor.fetchall()
    Ques_counter = 0

    if activity_type == 'multiple-choice' or activity_type == 'true-false':
        for i in submission_form:
            for j in i['options']:
                answer_query = "INSERT INTO Answers (question_id, answer, is_correct) Values (%s,'%s',%s)"\
                    % (Ques_ids[Ques_counter][0], j, i['correctAnswer'] == j)
                Cursor.execute(answer_query)
                Connection.commit()
            Ques_counter += 1
    else:
        for i in submission_form:
            answer_query = "INSERT INTO Answers (question_id, answer, is_correct) Values (%s,'%s',%s)"\
                % (Ques_ids[Ques_counter][0], i['correctAnswer'], 1)
            Cursor.execute(answer_query)
            Connection.commit()
            Ques_counter += 1

    Ques_counter = 0

    for i in submission_form:
        submitted_answer_query = "INSERT INTO Submitted_Answers (question_id,submission_id,submitted_answer) Values (%s,%s,'%s')"\
            % (Ques_ids[Ques_counter][0], session['Submission'], i['submittedAnswer'])
        Ques_counter += 1
        Cursor.execute(submitted_answer_query)
        Connection.commit()

    return ({"submission_id": session['Submission']})


@application.route('/api/documents/<int:document_id>/activity/<string:activity_type>/results/<int:submission_id>', methods=['GET'])
def get_submission(document_id, activity_type, submission_id):
    sub_SQL = "SELECT q.question_id, q.question, sa.submitted_answer, a.answer_id, a.answer, a.is_correct FROM user_docs.Questions AS q JOIN user_docs.Submitted_Answers AS sa ON q.question_id = sa.question_id JOIN user_docs.Answers AS a ON a.question_id = sa.question_id WHERE q.submission_id ='%s' " % submission_id

    try:
        Cursor.execute(sub_SQL)
        results = Cursor.fetchall()
    except:
        print("Database Query failed")
        return ("Database Connection failed")
    print("log_results:", results)

    return (results)


@application.route('/api/user-profile', methods=['GET'])
def get_submissions():
    get_user = "SELECT email FROM user_docs.Users WHERE user_id = '%s' " % session['ID']
    try:
        Cursor.execute(get_user)
        user = Cursor.fetchall()
    except:
        print("Database Query failed")
        return ("Database Connection failed")

    submissions_query = "SELECT DISTINCT sub.submission_id, sub.document_id, d.title, d.subtitle, q.question_type FROM user_docs.Submissions AS sub JOIN user_docs.Questions AS q ON sub.submission_id = q.submission_id JOIN user_docs.Users AS u ON sub.user_id = u.user_id JOIN user_docs.Documents as d ON sub.document_id = d.document_id WHERE sub.user_id = '%s' " % session[
        'ID']
    try:
        Cursor.execute(submissions_query)
        submissions = Cursor.fetchall()
    except:
        print("Database Query failed")
        return ("Database Connection failed")

    return ({"user": user, "submissions": submissions})


@application.route('/api/submission', methods=['POST'])
def submission():
    # submission = request.get_json()

    # get_ques_ids = "SELECT question_id FROM Questions WHERE submission_id = %s " % session[
    #     'Submission']
    # Cursor.execute(get_ques_ids)
    # Ques_ids = Cursor.fetchall()
    # Ques_counter = 0

    # for i in submission['answers']:
    #     answer_query = "INSERT INTO Submitted_Answers (question_id,submission_id,submitted_answer) Values (%s,%s,'%s')"\
    #         % (Ques_ids[0][Ques_counter], session['Submission'], i)
    #     Ques_counter += 1
    #     Cursor.execute(answer_query)
    #     Connection.commit()
    pass


@application.route('/api/documents/<int:document_id>/activity/<string:activity_type>/results/<int:submission_id>/feedback_report', methods=['POST'])
def send_feedback_report(document_id, activity_type, submission_id):
    # receiving from POST: email / feedback-message
    # send message report to email API with document_id, activity_type
    # respond with success message

    email = request.form['email']
    feedback = request.form['feedback']

    message = MIMEMultipart()
    # Replace with your Gmail address
    message['From'] = 'studysmartai.fau@gmail.com'
    message['To'] = 'studysmartai.fau@gmail.com'
    message['Subject'] = f"submission: userId={session['ID']}/{document_id}/activity/{activity_type}/results/{submission_id}"
    body = f"Submission: userId={session['ID']}/{document_id}/activity/{activity_type}/results/{submission_id}\n\n" \
        f"Email: {email}\n\n" \
        f"Feedback:\n{feedback}" \


    # Attach the body of the email
    message.attach(MIMEText(body, 'plain'))

    # Connect to the SMTP server (in this case, Gmail's SMTP server)
    try:
        with smtplib.SMTP('smtp.gmail.com', 587) as server:
            server.starttls()  # Enable TLS
            # Use your App Password
            server.login('studysmartai.fau@gmail.com', 'ijlqlsdrizgbhtxh')

            # Send the email
            server.sendmail('studysmartai.fau@gmail.com',
                            'studysmartai.fau@gmail.com', message.as_string())
    except:
        print("Email Send Failed")
        return ("Email Send Failed")

    return ("Success")


if __name__ == "__main__":
    application.run(debug=True)
